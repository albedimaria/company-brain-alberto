"""Binary artifact generation: docx / pptx / pdf / xlsx.

Files are written under static/files/ and served by the backend at /files/.
The agent calls `generate_artifact(...)`; it returns an absolute artifact_url.
HTML/markdown decks are NOT handled here - those go inline in `answer`.
"""

from __future__ import annotations

import re
import uuid
from pathlib import Path
from typing import Any

import config

_FILES = Path(__file__).resolve().parent / "static" / "files"
_FILES.mkdir(parents=True, exist_ok=True)

_SUPPORTED = ("pdf", "docx", "pptx", "xlsx")


def _safe_name(name: str, fmt: str) -> str:
    stem = (
        re.sub(r"[^a-zA-Z0-9_-]+", "-", (name or "artifact").rsplit(".", 1)[0]).strip(
            "-"
        )
        or "artifact"
    )
    return f"{stem}-{uuid.uuid4().hex[:8]}.{fmt}"


def _latin1(text: str) -> str:
    """fpdf2 core fonts are latin-1; map the few unicode chars we emit."""
    repl = {
        "\u20ac": "EUR",
        "\u2013": "-",
        "\u2014": "-",
        "\u2018": "'",
        "\u2019": "'",
        "\u201c": '"',
        "\u201d": '"',
        "\u2022": "-",
        "\u00a0": " ",
    }
    for a, b in repl.items():
        text = text.replace(a, b)
    return text.encode("latin-1", "replace").decode("latin-1")


def _pdf(path: Path, title: str, sections: list[dict[str, Any]]) -> None:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    w = pdf.epw  # effective page width (inside margins)
    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(w, 10, _latin1(title))
    pdf.ln(2)
    for sec in sections:
        if sec.get("heading"):
            pdf.set_font("Helvetica", "B", 13)
            pdf.multi_cell(w, 8, _latin1(str(sec["heading"])))
        if sec.get("body"):
            pdf.set_font("Helvetica", "", 11)
            pdf.multi_cell(w, 6, _latin1(str(sec["body"])))
        pdf.ln(2)
    pdf.output(str(path))


def _docx(path: Path, title: str, sections: list[dict[str, Any]]) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)
    for sec in sections:
        if sec.get("heading"):
            doc.add_heading(str(sec["heading"]), level=1)
        if sec.get("body"):
            for line in str(sec["body"]).split("\n"):
                doc.add_paragraph(line)
    doc.save(str(path))


def _pptx(path: Path, title: str, sections: list[dict[str, Any]]) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    for sec in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(sec.get("heading", title))
        body = slide.placeholders[1].text_frame
        body.word_wrap = True
        lines = str(sec.get("body", "")).split("\n")
        body.text = lines[0] if lines else ""
        for line in lines[1:]:
            p = body.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
    prs.save(str(path))


def _xlsx(
    path: Path,
    title: str,
    sections: list[dict[str, Any]],
    table: list[list[Any]] | None,
) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append([title])
    ws.append([])
    if table:
        for row in table:
            ws.append(list(row))
    else:
        for sec in sections:
            if sec.get("heading"):
                ws.append([str(sec["heading"])])
            if sec.get("body"):
                for line in str(sec["body"]).split("\n"):
                    ws.append([line])
            ws.append([])
    wb.save(str(path))


def generate_artifact(
    format: str,
    title: str,
    sections: list[dict[str, Any]] | None = None,
    table: list[list[Any]] | None = None,
    filename: str | None = None,
) -> dict[str, Any]:
    """Create a binary artifact and return {'artifact_url': <absolute url>}."""
    fmt = (format or "").lower().lstrip(".")
    if fmt not in _SUPPORTED:
        return {"error": f"unsupported format '{format}'. Use one of {_SUPPORTED}."}
    sections = sections or []
    name = _safe_name(filename or title, fmt)
    path = _FILES / name
    if fmt == "pdf":
        _pdf(path, title, sections)
    elif fmt == "docx":
        _docx(path, title, sections)
    elif fmt == "pptx":
        _pptx(path, title, sections)
    elif fmt == "xlsx":
        _xlsx(path, title, sections, table)
    return {
        "artifact_url": f"{config.PUBLIC_BASE_URL}/files/{name}",
        "format": fmt,
        "filename": name,
    }
